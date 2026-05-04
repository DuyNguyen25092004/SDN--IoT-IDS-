#!/usr/bin/env python3
"""
make_slides.py  —  Generate minimalist B&W presentation for SDN IoT IDS project.

Design: clean white background, black typography, zero color noise.
Run:    python3 make_slides.py
Output: IDS_Presentation.pptx  (same directory)
"""

import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR_TYPE

# ──────────────────────────────────────────────────────────────────────────────
# Paths
# ──────────────────────────────────────────────────────────────────────────────
BASE = Path(__file__).resolve().parent
P1   = BASE / "Part1"
P2   = BASE / "Part2"
P3   = BASE / "part3"

IMG = {
    "class_dist":     P2 / "class_distribution.png",
    "conf_matrix":    P2 / "confusion_matrix_best.png",
    "feat_imp":       P2 / "feature_importance_rf.png",
    "model_cmp":      P2 / "model_comparison.png",
    "per_class":      P2 / "per_class_performance.png",
    "roc":            BASE / "roc_curves.png",
    "conf_p3":        P3 / "confusion_matrix.png",
    "latency":        P3 / "latency_boxplot.png",
    "metrics_bar":    P3 / "metrics_bar.png",
}

OUT = BASE / "IDS_Presentation.pptx"

# ──────────────────────────────────────────────────────────────────────────────
# Color palette  (pure B&W minimalism)
# ──────────────────────────────────────────────────────────────────────────────
BLACK      = RGBColor(0x00, 0x00, 0x00)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK       = RGBColor(0x1A, 0x1A, 0x1A)   # near-black for body text
MIDGRAY    = RGBColor(0x77, 0x77, 0x77)   # captions / secondary
LIGHTGRAY  = RGBColor(0xEE, 0xEE, 0xEE)   # divider lines / backgrounds
VERYLGHT   = RGBColor(0xF7, 0xF7, 0xF7)   # panel fill

# ──────────────────────────────────────────────────────────────────────────────
# Helpers
# ──────────────────────────────────────────────────────────────────────────────
SW = Inches(13.33)   # slide width
SH = Inches(7.5)     # slide height

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs


def blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank_layout)


def bg_fill(slide, color: RGBColor = WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_w=Pt(0.75)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, l, t, w, h)   # 1 = MSO_SHAPE_TYPE.RECTANGLE
    shape.line.width = line_w
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=Pt(18), bold=False, color=DARK,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = "Calibri"
    return txb


def add_para(tf, text, size=Pt(16), bold=False, color=DARK,
             align=PP_ALIGN.LEFT, space_before=Pt(4), italic=False):
    """Append a paragraph to an existing text frame."""
    from pptx.util import Pt as _Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text   = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = "Calibri"
    return p


def hline(slide, t, l=Inches(0.6), w=None, color=BLACK, thickness=Pt(1)):
    """Horizontal rule."""
    if w is None:
        w = SW - Inches(1.2)
    ln = slide.shapes.add_shape(1, l, t, w, Pt(1))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()


def add_image(slide, path, l, t, w, h=None, max_h=None):
    """Insert image. Preserves aspect ratio; caps height at max_h if given."""
    from PIL import Image as _PilImg
    p = Path(path)
    if not p.exists():
        ph = h or max_h or Inches(3)
        add_rect(slide, l, t, w, ph, fill_color=LIGHTGRAY, line_color=MIDGRAY)
        add_text(slide, f"[{p.name}]", l, t + Inches(0.1), w, Inches(0.5),
                 size=Pt(11), color=MIDGRAY, align=PP_ALIGN.CENTER)
        return
    if h:
        slide.shapes.add_picture(str(path), l, t, w, h)
        return
    # Compute aspect-ratio-correct height and cap at max_h
    iw, ih = _PilImg.open(str(p)).size
    nat_h  = int(w * ih / iw)
    if max_h and nat_h > max_h:
        # Scale by height, center horizontally
        nh  = max_h
        nw  = int(max_h * iw / ih)
        off = max(0, (w - nw) // 2)
        slide.shapes.add_picture(str(path), l + off, t, nw, nh)
    else:
        slide.shapes.add_picture(str(path), l, t, w, nat_h)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    """Slide 1 — Cover"""
    s = blank_slide(prs)
    bg_fill(s, WHITE)

    # thick left accent bar
    add_rect(s, Inches(0), Inches(0), Inches(0.35), SH, fill_color=BLACK)

    # main title
    add_text(s,
             "Hệ thống phát hiện và chặn\nxâm nhập MQTT trên nền tảng SDN",
             Inches(0.7), Inches(1.8), Inches(8.5), Inches(2.4),
             size=Pt(36), bold=True, color=BLACK, align=PP_ALIGN.LEFT)

    # subtitle
    add_text(s,
             "IDS + XGBoost + Ryu Controller + Mininet",
             Inches(0.7), Inches(4.4), Inches(8), Inches(0.6),
             size=Pt(18), color=MIDGRAY, align=PP_ALIGN.LEFT)

    hline(s, Inches(4.15), l=Inches(0.7), w=Inches(7))

    # team / date block (bottom right)
    add_text(s,
             "Báo cáo cuối kỳ  ·  SDN  ·  2026",
             Inches(0.7), Inches(5.8), Inches(8), Inches(0.5),
             size=Pt(13), color=MIDGRAY, italic=True)

    # right column — system keyword tags
    tags = ["MQTT", "SDN / Ryu", "XGBoost", "Mininet", "IDS"]
    for i, tag in enumerate(tags):
        add_text(s, tag,
                 Inches(11.2), Inches(1.5 + i * 0.85), Inches(1.8), Inches(0.65),
                 size=Pt(14), bold=True, color=BLACK, align=PP_ALIGN.CENTER)
        add_rect(s, Inches(11.1), Inches(1.45 + i * 0.85),
                 Inches(2.0), Inches(0.65), line_color=BLACK)


def slide_section(prs, number: str, title: str, subtitle: str = ""):
    """Section divider — black background"""
    s = blank_slide(prs)
    bg_fill(s, BLACK)

    add_text(s, number,
             Inches(1), Inches(1.5), Inches(2), Inches(1),
             size=Pt(72), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    hline(s, Inches(2.8), l=Inches(1), w=Inches(11.3), color=WHITE, thickness=Pt(0.75))

    add_text(s, title,
             Inches(1), Inches(3.0), Inches(11), Inches(1.5),
             size=Pt(34), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    if subtitle:
        add_text(s, subtitle,
                 Inches(1), Inches(4.6), Inches(11), Inches(0.8),
                 size=Pt(17), color=RGBColor(0xAA, 0xAA, 0xAA), italic=True)


def slide_content(prs, title: str, bullets: list, footnote: str = ""):
    """Standard content slide with bullet list."""
    s = blank_slide(prs)
    bg_fill(s, WHITE)

    # title
    add_text(s, title,
             Inches(0.65), Inches(0.38), Inches(12), Inches(0.8),
             size=Pt(26), bold=True, color=BLACK, align=PP_ALIGN.LEFT)
    hline(s, Inches(1.22), l=Inches(0.65), color=BLACK, thickness=Pt(1.2))

    # body text box
    txb = slide.shapes.add_textbox(Inches(0.65), Inches(1.4),
                                    Inches(12.0), Inches(5.5)) \
          if False else None  # we use manual approach below

    y = Inches(1.4)
    for item in bullets:
        if isinstance(item, str):
            add_text(s, item,
                     Inches(0.85), y, Inches(11.6), Inches(0.55),
                     size=Pt(17), color=DARK)
            y += Inches(0.55)
        elif isinstance(item, tuple):
            lvl, text = item
            indent = Inches(0.85 + lvl * 0.4)
            bsz = Pt(17 - lvl * 1.5)
            bcol = DARK if lvl == 0 else MIDGRAY
            add_text(s, ("  · " if lvl > 0 else "— ") + text,
                     indent, y, Inches(11.8 - lvl * 0.4), Inches(0.5),
                     size=bsz, color=bcol)
            y += Inches(0.48)

    if footnote:
        add_text(s, footnote,
                 Inches(0.65), Inches(6.9), Inches(12), Inches(0.45),
                 size=Pt(11), color=MIDGRAY, italic=True)
    return s


def build_content_slide(prs, title: str, items: list, footnote: str = ""):
    """Build a content slide from a list of (level, text) or plain strings."""
    s = blank_slide(prs)
    bg_fill(s, WHITE)

    add_text(s, title,
             Inches(0.65), Inches(0.35), Inches(12), Inches(0.85),
             size=Pt(26), bold=True, color=BLACK)
    hline(s, Inches(1.25), l=Inches(0.65), color=BLACK, thickness=Pt(1.5))

    y = Inches(1.45)
    for item in items:
        if isinstance(item, str):
            lvl, text = 0, item
        else:
            lvl, text = item
        prefix = "— " if lvl == 0 else "  · "
        indent = Inches(0.85 + lvl * 0.45)
        fsz    = Pt(16 - lvl * 1)
        col    = DARK if lvl == 0 else MIDGRAY
        add_text(s, prefix + text,
                 indent, y, Inches(11.8 - lvl * 0.45), Inches(0.46),
                 size=fsz, color=col)
        y += Inches(0.48)

    if footnote:
        add_text(s, footnote,
                 Inches(0.65), Inches(6.9), Inches(12), Inches(0.45),
                 size=Pt(11), color=MIDGRAY, italic=True)
    return s


def build_image_slide(prs, title: str, img_path, caption: str = "",
                      img_left=None, img_top=None,
                      img_w=Inches(11), img_h=None):
    s = blank_slide(prs)
    bg_fill(s, WHITE)

    add_text(s, title,
             Inches(0.65), Inches(0.35), Inches(12), Inches(0.85),
             size=Pt(26), bold=True, color=BLACK)
    hline(s, Inches(1.25), l=Inches(0.65), color=BLACK, thickness=Pt(1.5))

    l = img_left if img_left is not None else Inches(1.15)
    t = img_top  if img_top  is not None else Inches(1.55)

    # Cap image height so it never reaches the caption / slide bottom
    cap_gap = Inches(0.55) if caption else Inches(0.15)
    max_h   = Inches(7.05) - t - cap_gap

    add_image(s, img_path, l, t, img_w, img_h, max_h=max_h)

    if caption:
        add_text(s, caption,
                 Inches(0.65), Inches(6.9), Inches(12), Inches(0.45),
                 size=Pt(12), color=MIDGRAY, italic=True,
                 align=PP_ALIGN.CENTER)
    return s


def build_two_image_slide(prs, title: str,
                          img1, img2,
                          cap1="", cap2="",
                          footnote=""):
    s = blank_slide(prs)
    bg_fill(s, WHITE)

    add_text(s, title,
             Inches(0.65), Inches(0.35), Inches(12), Inches(0.85),
             size=Pt(26), bold=True, color=BLACK)
    hline(s, Inches(1.25), l=Inches(0.65), color=BLACK, thickness=Pt(1.5))

    half_w = Inches(5.8)
    max_h  = Inches(4.9)   # safe ceiling: 1.55 + 4.9 + 0.45 caption < 7.0"
    add_image(s, img1, Inches(0.4), Inches(1.55), half_w, max_h=max_h)
    add_image(s, img2, Inches(6.9), Inches(1.55), half_w, max_h=max_h)

    if cap1:
        add_text(s, cap1, Inches(0.4), Inches(6.45), half_w, Inches(0.4),
                 size=Pt(11), color=MIDGRAY, italic=True, align=PP_ALIGN.CENTER)
    if cap2:
        add_text(s, cap2, Inches(6.9), Inches(6.45), half_w, Inches(0.4),
                 size=Pt(11), color=MIDGRAY, italic=True, align=PP_ALIGN.CENTER)
    if footnote:
        add_text(s, footnote,
                 Inches(0.65), Inches(6.9), Inches(12), Inches(0.45),
                 size=Pt(11), color=MIDGRAY, italic=True)
    return s


def build_stats_slide(prs, title: str, stats: list, note: str = ""):
    """
    stats: list of (label, value, unit) tuples.
    Lays them out as large centered numbers.
    """
    s = blank_slide(prs)
    bg_fill(s, WHITE)

    add_text(s, title,
             Inches(0.65), Inches(0.35), Inches(12), Inches(0.85),
             size=Pt(26), bold=True, color=BLACK)
    hline(s, Inches(1.25), l=Inches(0.65), color=BLACK, thickness=Pt(1.5))

    n    = len(stats)
    col_w = (SW - Inches(1.3)) / n
    for i, (label, value, unit) in enumerate(stats):
        x = Inches(0.65) + i * col_w
        # big number
        add_text(s, str(value),
                 x, Inches(2.3), col_w, Inches(1.8),
                 size=Pt(54), bold=True, color=BLACK, align=PP_ALIGN.CENTER)
        # unit
        add_text(s, unit,
                 x, Inches(4.15), col_w, Inches(0.55),
                 size=Pt(14), color=MIDGRAY, align=PP_ALIGN.CENTER)
        # label
        add_text(s, label,
                 x, Inches(4.75), col_w, Inches(0.6),
                 size=Pt(14), bold=True, color=DARK, align=PP_ALIGN.CENTER)
        # vertical rule between cells
        if i > 0:
            add_rect(s, x - Inches(0.02), Inches(2.0),
                     Pt(1.2), Inches(3.5), fill_color=LIGHTGRAY)

    if note:
        add_text(s, note, Inches(0.65), Inches(6.85),
                 Inches(12), Inches(0.45),
                 size=Pt(11), color=MIDGRAY, italic=True, align=PP_ALIGN.CENTER)
    return s


def build_flow_slide(prs, title: str, steps: list, footnote: str = ""):
    """
    Horizontal flow diagram.  steps = list of (label, detail).
    """
    s = blank_slide(prs)
    bg_fill(s, WHITE)

    add_text(s, title,
             Inches(0.65), Inches(0.35), Inches(12), Inches(0.85),
             size=Pt(26), bold=True, color=BLACK)
    hline(s, Inches(1.25), l=Inches(0.65), color=BLACK, thickness=Pt(1.5))

    n      = len(steps)
    box_w  = Inches(1.7)
    box_h  = Inches(1.3)
    gap    = Inches(0.22)
    total  = n * box_w + (n - 1) * gap
    start  = (SW - total) / 2
    top    = Inches(2.8)

    for i, (label, detail) in enumerate(steps):
        x = start + i * (box_w + gap)
        # box
        add_rect(s, x, top, box_w, box_h,
                 fill_color=WHITE, line_color=BLACK, line_w=Pt(1.5))
        # label
        add_text(s, label,
                 x + Inches(0.05), top + Inches(0.12),
                 box_w - Inches(0.1), Inches(0.55),
                 size=Pt(13), bold=True, color=BLACK, align=PP_ALIGN.CENTER)
        # detail
        add_text(s, detail,
                 x + Inches(0.05), top + Inches(0.7),
                 box_w - Inches(0.1), Inches(0.55),
                 size=Pt(10), color=MIDGRAY, align=PP_ALIGN.CENTER)
        # arrow (except last)
        if i < n - 1:
            ax = x + box_w + Inches(0.03)
            add_text(s, "→",
                     ax, top + Inches(0.4),
                     gap + Inches(0.15), Inches(0.5),
                     size=Pt(18), bold=True, color=BLACK, align=PP_ALIGN.CENTER)

    # detail row
    y_det = top + box_h + Inches(0.55)
    for i, (label, detail) in enumerate(steps):
        x = start + i * (box_w + gap)
        add_text(s, detail,
                 x - Inches(0.05), y_det,
                 box_w + Inches(0.1), Inches(1.2),
                 size=Pt(11), color=DARK, align=PP_ALIGN.CENTER)

    if footnote:
        add_text(s, footnote,
                 Inches(0.65), Inches(6.9), Inches(12), Inches(0.45),
                 size=Pt(11), color=MIDGRAY, italic=True)
    return s


def build_topology_slide(prs):
    """Clean star topology: publishers left, switch center, subscribers/broker/attacker."""
    s = blank_slide(prs)
    bg_fill(s, WHITE)

    add_text(s, "Topology Mininet",
             Inches(0.65), Inches(0.35), Inches(12), Inches(0.85),
             size=Pt(26), bold=True, color=BLACK)
    hline(s, Inches(1.25), l=Inches(0.65), color=BLACK, thickness=Pt(1.5))

    NW, NH = Inches(1.3), Inches(0.55)   # node width / height

    # ── s1 Switch (center) ─────────────────────────────────────────────────
    SW_L, SW_T = Inches(4.6), Inches(3.15)
    SW_W, SW_H = Inches(1.9),  Inches(0.75)
    SW_CX = SW_L + SW_W / 2   # 5.55"
    SW_CY = SW_T + SW_H / 2   # 3.525"

    add_rect(s, SW_L, SW_T, SW_W, SW_H, fill_color=BLACK)
    add_text(s, "s1  (OVS)", SW_L, SW_T, SW_W, SW_H,
             size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ── Helpers ─────────────────────────────────────────────────────────────
    def nd(lx, ty, lbl, ip, fill=WHITE, tc=BLACK):
        add_rect(s, lx, ty, NW, NH, fill_color=fill, line_color=BLACK, line_w=Pt(1.2))
        add_text(s, f"{lbl}\n{ip}", lx, ty, NW, NH,
                 size=Pt(8.5), color=tc, align=PP_ALIGN.CENTER)
        return lx + NW / 2, ty + NH / 2

    def wire(x1, y1, x2, y2, lw=Pt(1.0)):
        c = s.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                                   int(x1), int(y1), int(x2), int(y2))
        c.line.color.rgb = BLACK
        c.line.width = lw

    # ── Publishers h1–h6 (left column) ───────────────────────────────────────
    PUB_X = Inches(0.3)
    add_text(s, "Publishers", PUB_X, Inches(1.32), NW, Inches(0.24),
             size=Pt(8), bold=True, color=MIDGRAY, align=PP_ALIGN.CENTER)
    pub_data = [("h1","10.0.0.1"),("h2","10.0.0.2"),("h3","10.0.0.3"),
                ("h4","10.0.0.4"),("h5","10.0.0.5"),("h6","10.0.0.6")]
    for i, (lbl, ip) in enumerate(pub_data):
        ny = Inches(1.55 + i * 0.7)
        _, ncy = nd(PUB_X, ny, lbl, ip)
        # each node right-edge → switch left-edge (fan converges at SW_CY)
        wire(PUB_X + NW, ncy, SW_L, SW_CY)
    # group label mid-way between col and switch
    add_text(s, "eth1–6", Inches(2.4), Inches(3.42), Inches(0.8), Inches(0.26),
             size=Pt(7.5), color=MIDGRAY, italic=True, align=PP_ALIGN.CENTER)

    # ── Subscribers h7–h8 (right of switch) ─────────────────────────────────
    SUB_X = Inches(7.2)
    add_text(s, "Subscribers", SUB_X, Inches(2.45), NW, Inches(0.24),
             size=Pt(8), bold=True, color=MIDGRAY, align=PP_ALIGN.CENTER)
    sub_data = [("h7","10.0.0.7"),("h8","10.0.0.8")]
    for i, (lbl, ip) in enumerate(sub_data):
        ny = Inches(2.72 + i * 0.95)
        _, ncy = nd(SUB_X, ny, lbl, ip)
        # switch right-edge → node left-edge
        wire(SW_L + SW_W, SW_CY, SUB_X, ncy)
    add_text(s, "eth7–8", Inches(6.54), Inches(3.42), Inches(0.62), Inches(0.26),
             size=Pt(7.5), color=MIDGRAY, italic=True, align=PP_ALIGN.CENTER)

    # ── hbroker (above switch, dark fill) ────────────────────────────────────
    BR_L = SW_CX - NW / 2   # horizontally centered on switch
    BR_T = Inches(1.42)
    brcx, _ = nd(BR_L, BR_T, "hbroker", "10.0.0.10", fill=DARK, tc=WHITE)
    wire(brcx, BR_T + NH, SW_CX, SW_T)   # bottom of broker → top of switch
    add_text(s, "eth9", SW_CX + Inches(0.1), Inches(2.56), Inches(0.4), Inches(0.24),
             size=Pt(7.5), color=MIDGRAY, italic=True)

    # ── hattacker (below switch) ────────────────────────────────────────────
    AT_L = SW_CX - NW / 2
    AT_T = Inches(5.22)
    atcx, _ = nd(AT_L, AT_T, "hattacker", "10.0.0.99")
    wire(atcx, AT_T, SW_CX, SW_T + SW_H)  # top of attacker → bottom of switch
    add_text(s, "eth10", SW_CX + Inches(0.1), Inches(4.64), Inches(0.46), Inches(0.24),
             size=Pt(7.5), color=MIDGRAY, italic=True)

    # ── eth11 mirror note (below subscribers, no host) ─────────────────────────
    add_text(s, "s1-eth11  (mirror port, no host)\ntshark listens → IDS API",
             Inches(7.2), Inches(4.9), Inches(3.5), Inches(0.55),
             size=Pt(8.5), color=MIDGRAY, italic=True)

    # ── Bottom info bar (Ryu / IDS / tshark / broker context) ──────────────────
    add_rect(s, Inches(0.3), Inches(6.32), Inches(12.73), Inches(0.65),
             fill_color=VERYLGHT, line_color=LIGHTGRAY, line_w=Pt(0.5))
    add_text(s,
             "Ryu Controller  :6633 / :8080  (OpenFlow 1.3)  ·  "
             "IDS API  :5000  (XGBoost v5)  ·  "
             "tshark on s1-eth11  ·  Mosquitto  :1883 on hbroker",
             Inches(0.5), Inches(6.36), Inches(12.33), Inches(0.56),
             size=Pt(9.5), color=DARK, align=PP_ALIGN.CENTER)
    return s


def build_attack_single_slide(prs):
    s = blank_slide(prs)
    bg_fill(s, WHITE)

    add_text(s, "Kịch bản 1 — Tấn công đơn lẻ (tuần tự)",
             Inches(0.65), Inches(0.35), Inches(12), Inches(0.85),
             size=Pt(26), bold=True, color=BLACK)
    hline(s, Inches(1.25), l=Inches(0.65), color=BLACK, thickness=Pt(1.5))

    attacks = [
        ("A1", "MQTT Flood",     "attack1_mqtt_flood.py\n300 msg/s, 5 threads\nDoS — broker saturation"),
        ("A2", "DoS (TCP)",       "attack2_dos.py\nTCP connection flood\nExhaust broker sockets"),
        ("A3", "Brute Force",    "attack3_brute_force.py\n300 user/pass pairs\nCONNECT spam"),
        ("A4", "Malformed",      "attack4_malformed.py\nBad MQTT headers\nParser exploit attempt"),
        ("A5", "Slow Drip",      "attack5_slow_drip.py\n1.5 msg/s, low-noise\nStealth exfiltration"),
    ]

    bw = Inches(2.15)
    bh = Inches(3.5)
    gap = Inches(0.2)
    total = len(attacks) * bw + (len(attacks) - 1) * gap
    sx = (SW - total) / 2

    for i, (tag, name, detail) in enumerate(attacks):
        x = sx + i * (bw + gap)
        # card
        add_rect(s, x, Inches(1.65), bw, bh, fill_color=WHITE, line_color=BLACK, line_w=Pt(1.5))
        # tag
        add_rect(s, x, Inches(1.65), bw, Inches(0.52), fill_color=BLACK)
        add_text(s, tag, x, Inches(1.65), bw, Inches(0.52),
                 size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # name
        add_text(s, name, x + Inches(0.1), Inches(2.25),
                 bw - Inches(0.2), Inches(0.7),
                 size=Pt(14), bold=True, color=BLACK, align=PP_ALIGN.CENTER)
        # detail
        add_text(s, detail, x + Inches(0.12), Inches(3.0),
                 bw - Inches(0.24), Inches(2.0),
                 size=Pt(11), color=DARK, align=PP_ALIGN.CENTER)

    add_text(s,
             "Mỗi kịch bản: 1 attacker (hattacker) → 1 broker (10.0.0.10)  ·  Reset IDS state trước mỗi lần chạy",
             Inches(0.65), Inches(5.35), Inches(12.0), Inches(0.42),
             size=Pt(12), color=MIDGRAY, italic=True, align=PP_ALIGN.CENTER)
    return s


def build_attack_killchain_slide(prs):
    s = blank_slide(prs)
    bg_fill(s, WHITE)

    add_text(s, "Kịch bản 2 — Kill-chain (1 kẻ tấn công, đa giai đoạn)",
             Inches(0.65), Inches(0.35), Inches(12.3), Inches(0.85),
             size=Pt(24), bold=True, color=BLACK)
    hline(s, Inches(1.25), l=Inches(0.65), color=BLACK, thickness=Pt(1.5))

    add_text(s, "Script: combined_attack.py  ·  Profile: fast (~60s)  ·  Kẻ tấn công: hattacker (10.0.0.99)",
             Inches(0.65), Inches(1.35), Inches(12), Inches(0.4),
             size=Pt(12), color=MIDGRAY, italic=True)

    phases = [
        ("t=0s",  "A  RECON",         "Malformed probe\n10 pkt/s · 5s\nFingerprint broker"),
        ("t=8s",  "B  CREDENTIAL",    "Brute-force CONNECT\ndelay=0.05s · 12s\nThu thập credential"),
        ("t=22s", "C  IMPACT",        "MQTT Flood\n300 msg/s · 15s\nSmokescreen DoS"),
        ("t=40s", "D  EXFIL",         "Slow-drip publish\n1.5 msg/s · 20s\nStealth data leak"),
    ]

    pw  = Inches(2.6)
    ph  = Inches(3.2)
    gap = Inches(0.35)
    total = len(phases) * pw + (len(phases) - 1) * gap
    px  = (SW - total) / 2
    ty  = Inches(1.9)

    for i, (ts, phase_name, desc) in enumerate(phases):
        x = px + i * (pw + gap)
        # timeline marker
        add_text(s, ts,
                 x, ty - Inches(0.45), pw, Inches(0.38),
                 size=Pt(11), bold=True, color=MIDGRAY, align=PP_ALIGN.CENTER)
        # card outline
        add_rect(s, x, ty, pw, ph, fill_color=WHITE, line_color=BLACK, line_w=Pt(1.5))
        # phase header bar
        add_rect(s, x, ty, pw, Inches(0.6), fill_color=BLACK)
        add_text(s, phase_name,
                 x + Inches(0.08), ty + Inches(0.05), pw - Inches(0.16), Inches(0.5),
                 size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # desc
        add_text(s, desc,
                 x + Inches(0.12), ty + Inches(0.7), pw - Inches(0.24), Inches(2.4),
                 size=Pt(12), color=DARK, align=PP_ALIGN.CENTER)
        # arrow between cards
        if i < len(phases) - 1:
            add_text(s, "→",
                     x + pw + Inches(0.05), ty + Inches(1.3), gap + Inches(0.2), Inches(0.55),
                     size=Pt(20), bold=True, color=BLACK, align=PP_ALIGN.CENTER)

    # expected IDS behaviour row
    add_text(s, "IDS kỳ vọng:",
             Inches(0.65), Inches(5.3), Inches(1.8), Inches(0.45),
             size=Pt(13), bold=True, color=BLACK)
    add_text(s,
             "Block kích hoạt ở pha A hoặc B (threat score vượt ngưỡng)  ·  "
             "Ryu cài DROP rule cho 10.0.0.99  ·  Pha C & D bị chặn tại data plane",
             Inches(2.5), Inches(5.3), Inches(10.2), Inches(0.7),
             size=Pt(13), color=DARK)
    return s


def build_attack_multi_slide(prs):
    s = blank_slide(prs)
    bg_fill(s, WHITE)

    add_text(s, "Kịch bản 3 — Đa kẻ tấn công (Multi-Attacker)",
             Inches(0.65), Inches(0.35), Inches(12), Inches(0.85),
             size=Pt(26), bold=True, color=BLACK)
    hline(s, Inches(1.25), l=Inches(0.65), color=BLACK, thickness=Pt(1.5))

    add_text(s, "Script: scenario_multi_attacker.py  (part3)",
             Inches(0.65), Inches(1.35), Inches(12), Inches(0.4),
             size=Pt(13), color=MIDGRAY, italic=True)

    # 2-column layout: Attacker 1 and Attacker 2
    col_w = Inches(5.4)
    for col, (host, ip, atks) in enumerate([
        ("h_attacker", "10.0.0.99",
         ["A3 Brute Force CONNECT", "A5 Slow Drip (stealth channel)"]),
        ("h5",         "10.0.0.5",
         ["A1 MQTT Flood (smokescreen)", "A4 Malformed packets"]),
    ]):
        x = Inches(0.65) + col * (col_w + Inches(1.2))
        add_rect(s, x, Inches(1.85), col_w, Inches(3.6),
                 fill_color=VERYLGHT, line_color=BLACK, line_w=Pt(1.2))
        add_text(s, f"{host}\n{ip}",
                 x + Inches(0.15), Inches(1.95), col_w - Inches(0.3), Inches(0.75),
                 size=Pt(14), bold=True, color=BLACK, align=PP_ALIGN.CENTER)
        hline(s, Inches(2.75), l=x + Inches(0.2), w=col_w - Inches(0.4),
              color=MIDGRAY, thickness=Pt(0.5))
        for j, atk in enumerate(atks):
            add_text(s, "· " + atk,
                     x + Inches(0.25), Inches(2.9 + j * 0.65),
                     col_w - Inches(0.5), Inches(0.55),
                     size=Pt(13), color=DARK)

    # broker target
    add_text(s, "hbroker\n10.0.0.10\n(MQTT Broker)",
             Inches(5.6), Inches(2.7), Inches(2.0), Inches(1.4),
             size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(5.55), Inches(2.65), Inches(2.1), Inches(1.5),
             fill_color=BLACK)
    add_text(s, "hbroker\n10.0.0.10\n(MQTT Broker)",
             Inches(5.55), Inches(2.65), Inches(2.1), Inches(1.5),
             size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    rows = [
        "IDS xử lý đồng thời 2 IP nguồn — threat score độc lập cho mỗi IP.",
        "Ryu cài DROP rule riêng biệt khi từng IP vượt ngưỡng.",
        "Kịch bản kiểm tra khả năng xử lý concurrent attack trên cùng 1 switch.",
    ]
    for j, row in enumerate(rows):
        add_text(s, "· " + row,
                 Inches(0.65), Inches(5.65 + j * 0.5), Inches(12), Inches(0.45),
                 size=Pt(13), color=DARK)

    return s


def build_arch_slide(prs):
    """System architecture flow."""
    s = blank_slide(prs)
    bg_fill(s, WHITE)

    add_text(s, "Kiến trúc hệ thống",
             Inches(0.65), Inches(0.35), Inches(12), Inches(0.85),
             size=Pt(26), bold=True, color=BLACK)
    hline(s, Inches(1.25), l=Inches(0.65), color=BLACK, thickness=Pt(1.5))

    components = [
        ("IoT Hosts\n(h1–h8)",     "Publishers &\nSubscribers\nMQTT over TCP"),
        ("s1\nOVS Switch",         "OpenFlow 1.3\nMirror → eth11\nDrop rules"),
        ("tshark\nCapture",        "traffic_capture.py\n12 raw features\nPOST /predict"),
        ("IDS API\nFlask :5000",   "XGBoost v5\n16 features\nThreat score"),
        ("Ryu\nController :6633",  "L2 learning\nREST block\nDROP rule"),
    ]

    bw  = Inches(2.0)
    bh  = Inches(2.0)
    gap = Inches(0.22)
    total = len(components) * bw + (len(components) - 1) * gap
    sx  = (SW - total) / 2
    ty  = Inches(2.3)

    for i, (name, detail) in enumerate(components):
        x = sx + i * (bw + gap)
        add_rect(s, x, ty, bw, bh, fill_color=WHITE, line_color=BLACK, line_w=Pt(1.5))
        add_text(s, name,
                 x + Inches(0.08), ty + Inches(0.12),
                 bw - Inches(0.16), Inches(0.7),
                 size=Pt(12), bold=True, color=BLACK, align=PP_ALIGN.CENTER)
        hline(s, ty + Inches(0.85), l=x + Inches(0.12),
              w=bw - Inches(0.24), color=LIGHTGRAY, thickness=Pt(0.5))
        add_text(s, detail,
                 x + Inches(0.08), ty + Inches(0.92),
                 bw - Inches(0.16), Inches(1.0),
                 size=Pt(10), color=MIDGRAY, align=PP_ALIGN.CENTER)
        if i < len(components) - 1:
            add_text(s, "→",
                     x + bw + Inches(0.02), ty + Inches(0.75),
                     gap + Inches(0.16), Inches(0.5),
                     size=Pt(18), bold=True, color=BLACK, align=PP_ALIGN.CENTER)

    # MQTT broker below switch
    add_text(s, "hbroker\n(Mosquitto :1883)",
             Inches(4.35), Inches(5.0), Inches(2.4), Inches(0.8),
             size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(4.3), Inches(4.95), Inches(2.5), Inches(0.85),
             fill_color=BLACK)
    add_text(s, "hbroker\n(Mosquitto :1883)",
             Inches(4.3), Inches(4.95), Inches(2.5), Inches(0.85),
             size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s, "↕ MQTT :1883",
             Inches(4.7), Inches(4.45), Inches(1.8), Inches(0.45),
             size=Pt(10), color=MIDGRAY, italic=True, align=PP_ALIGN.CENTER)

    return s


def build_thank_you(prs):
    s = blank_slide(prs)
    bg_fill(s, BLACK)

    add_text(s, "Cảm ơn đã lắng nghe.",
             Inches(1), Inches(2.4), Inches(11), Inches(1.5),
             size=Pt(42), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    hline(s, Inches(4.1), l=Inches(1), w=Inches(6), color=WHITE, thickness=Pt(0.75))

    add_text(s,
             "Source code: github.com/SDN--IoT-IDS-\n"
             "Dataset: MQTTset (Kaggle)  ·  Model: XGBoost v5  ·  Controller: Ryu 4.x",
             Inches(1), Inches(4.4), Inches(11), Inches(1.2),
             size=Pt(15), color=RGBColor(0xAA, 0xAA, 0xAA), italic=True)
    return s


# ══════════════════════════════════════════════════════════════════════════════
#  MAIN — assemble all slides
# ══════════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()

    # ─── 1. Cover ─────────────────────────────────────────────────────────────
    slide_title(prs)

    # ─── 2. Outline ───────────────────────────────────────────────────────────
    build_content_slide(prs, "Nội dung báo cáo", [
        (0, "Đặt vấn đề — MQTT, SDN và bài toán IDS"),
        (0, "Mô hình AI — XGBoost, MQTTset, 16 đặc trưng"),
        (0, "Kết quả huấn luyện — Accuracy, F1, Confusion Matrix, ROC"),
        (0, "Triển khai hệ thống — IDS API, Traffic Capture, Ryu Controller"),
        (0, "Topology Mininet — 8 hosts, 1 switch, 1 broker"),
        (0, "Kịch bản tấn công — Đơn lẻ, Kill-chain, Đa kẻ tấn công"),
        (0, "Đánh giá hệ thống — Per-class metrics, Latency"),
        (0, "Kết luận & Hướng phát triển"),
    ])

    # ─── SECTION 1: Introduction ──────────────────────────────────────────────
    slide_section(prs, "01", "Đặt vấn đề",
                  "MQTT · SDN · Bài toán xâm nhập IoT")

    build_content_slide(prs, "MQTT — Message Queuing Telemetry Transport", [
        (0, "Giao thức nhắn tin nhẹ (lightweight), thiết kế cho IoT / M2M"),
        (1, "Mô hình Publish / Subscribe qua broker trung gian"),
        (1, "Chạy trên TCP/IP, port mặc định 1883 (plain) / 8883 (TLS)"),
        (1, "Chuẩn OASIS — dùng rộng rãi trong smart home, industrial IoT"),
        (0, "Ưu điểm: nhỏ gọn (2-byte header), hỗ trợ QoS 0/1/2, retain, wildcard topics"),
        (0, "Nhược điểm bảo mật:"),
        (1, "Không có xác thực bắt buộc trong spec cũ (3.1 / 3.1.1)"),
        (1, "Broker dễ bị DoS, brute-force, malformed packet injection"),
        (1, "Không có encryption tích hợp — dữ liệu đi plain-text"),
    ], footnote="MQTT v3.1.1 (ISO/IEC 20922)  ·  MQTT v5.0 thêm enhanced auth nhưng chưa phổ biến")

    build_content_slide(prs, "SDN — Software-Defined Networking", [
        (0, "Tách rời control plane và data plane"),
        (1, "Control plane: bộ điều khiển tập trung (Ryu, OpenDaylight, ONOS)"),
        (1, "Data plane: thiết bị chuyển tiếp (OVS, switch OpenFlow)"),
        (0, "Lợi thế cho IDS:"),
        (1, "Cài DROP rule tức thì ở data plane → block ở wire speed"),
        (1, "Không cần cài agent trên từng host"),
        (1, "Tầm nhìn toàn cục — thấy tất cả lưu lượng qua switch"),
        (0, "Trong đề tài: Ryu controller + OVS s1 + OpenFlow 1.3"),
        (1, "REST API /ids/block — IDS gọi để cài DROP rule"),
        (1, "Port mirroring eth11 — tshark bắt traffic không ảnh hưởng luồng chính"),
    ])

    build_content_slide(prs, "Bài toán: Xâm nhập trong mạng IoT/MQTT", [
        (0, "5 loại tấn công cần phát hiện:"),
        (1, "Flood (DoS) — gửi ồ ạt PUBLISH làm broker quá tải"),
        (1, "DoS (TCP) — flood TCP connection để cạn socket broker"),
        (1, "Brute Force — thử hàng trăm cặp username/password CONNECT"),
        (1, "Malformed — gói MQTT sai header, khai thác lỗ hổng parser"),
        (1, "Slow Drip — exfiltration tốc độ thấp, ẩn trong traffic bình thường"),
        (0, "Thách thức:"),
        (1, "Traffic mix — normal & attack trên cùng luồng TCP port 1883"),
        (1, "Phân loại realtime — quyết định phải trong vài giây"),
        (1, "Block tự động — không có người vận hành can thiệp"),
    ])

    # ─── SECTION 2: AI Model ──────────────────────────────────────────────────
    slide_section(prs, "02", "Mô hình AI — XGBoost",
                  "MQTTset · 16 đặc trưng · 6 lớp phân loại")

    build_content_slide(prs, "Dataset — MQTTset", [
        (0, "Nguồn: Kaggle — MQTTset (Vaccari et al., 2020)"),
        (0, "Tổng: ~14.4M gói tin  ·  Train 8.4M / Test 6M"),
        (0, "6 lớp nhãn:"),
        (1, "legitimate  ·  bruteforce  ·  dos  ·  flood  ·  malformed  ·  slowite"),
        (0, "Mất cân bằng lớp: legitimate chiếm ~70% → áp dụng RandomUnderSampler"),
        (0, "33 feature gốc → chọn lọc 16 feature quan trọng nhất:"),
        (1, "12 per-packet: tcp.len, tcp.time_delta, tcp.flags, mqtt.msgtype, ..."),
        (1, "4 per-IP aggregate (cửa sổ trượt): time_delta_mean, time_delta_std, pkt_rate, pub_to_conn_ratio"),
    ], footnote="Feature lựa chọn dựa trên feature importance XGBoost và loại bỏ feature phụ thuộc nhau")

    build_image_slide(prs, "Phân bố lớp trong dataset",
                      IMG["class_dist"],
                      caption="Class distribution sau RandomUnderSampler — legitimate được hạ xuống cân bằng với attack classes",
                      img_left=Inches(1.5), img_top=Inches(1.55),
                      img_w=Inches(10))

    build_content_slide(prs, "Mô hình XGBoost — Cấu hình", [
        (0, "Thuật toán: XGBoost (Extreme Gradient Boosting) — multiclass softmax"),
        (0, "Hyperparameters tối ưu:"),
        (1, "n_estimators = 400  ·  best_iteration = 307  ·  max_depth = 8"),
        (1, "learning_rate = 0.1  ·  subsample = 0.8  ·  colsample_bytree = 0.8"),
        (1, "min_child_weight = 5  ·  gamma = 0.1  ·  reg_alpha / reg_lambda = 0.1 / 1.0"),
        (1, "tree_method = hist  ·  device = cuda  ·  early_stopping_rounds = 20"),
        (0, "Pipeline: StandardScaler → XGBClassifier → LabelEncoder (6 classes)"),
        (0, "Lý do chọn XGBoost:"),
        (1, "Xử lý tốt mixed-type feature (numeric + binary flags)"),
        (1, "Tốc độ inference nhanh (<1ms/sample) — phù hợp realtime IDS"),
        (1, "Khả năng giải thích cao qua feature importance"),
    ])

    build_stats_slide(prs, "Kết quả huấn luyện", [
        ("Validation\nAccuracy", "81.6%", ""),
        ("Test\nAccuracy",       "83.2%", ""),
        ("Test\nF1 (weighted)",  "83.8%", ""),
        ("Classes",              "6",     "loại phân loại"),
        ("Features",             "16",    "/ 33 đặc trưng"),
    ], note="Trained on 8.4M samples  ·  XGBoost v5  ·  Saved: 2026-05-01")

    build_two_image_slide(prs, "Confusion Matrix — Validation & Test",
                          IMG["conf_matrix"], IMG["conf_p3"],
                          cap1="Validation set confusion matrix",
                          cap2="Test set (deployment evaluation)")

    build_image_slide(prs, "Feature Importance",
                      IMG["feat_imp"],
                      caption="Top features: pkt_rate, time_delta_std, tcp.len, mqtt.msgtype — aggregate features đóng góp lớn nhất",
                      img_left=Inches(1.0), img_top=Inches(1.55),
                      img_w=Inches(11.3))

    build_image_slide(prs, "ROC Curves — 6 lớp phân loại",
                      IMG["roc"],
                      caption="AUC > 0.95 cho hầu hết các lớp. Slowite thấp nhất do đặc trưng chồng lấp với legitimate traffic",
                      img_left=Inches(1.5), img_top=Inches(1.55),
                      img_w=Inches(10))

    # ─── SECTION 3: Deployment ────────────────────────────────────────────────
    slide_section(prs, "03", "Triển khai hệ thống",
                  "IDS API · Traffic Capture · Ryu Controller")

    build_arch_slide(prs)

    build_content_slide(prs, "IDS API — Flask REST (ids_api.py)", [
        (0, "POST /predict — nhận 12 raw features + tính 4 aggregate per-IP"),
        (0, "Threat Score system (thay thế vote window):"),
        (1, "Mỗi IP có 1 điểm số (0–10), tăng khi detect attack, giảm theo exponential decay"),
        (1, "Nhiều loại attack cùng lúc đều cộng vào 1 pool — không reset cứng"),
        (1, "IDS_BLOCK_CONF=0.8 · IDS_VOTE_WINDOW=20 · IDS_AGG_WINDOW=10s"),
        (0, "Khi threat score vượt ngưỡng → gọi POST /ids/block lên Ryu"),
        (0, "Whitelist: broker (10.0.0.10), localhost — không bao giờ bị block"),
        (0, "GET /stats — xem số packet phân loại, block events, confidence"),
        (0, "POST /reset — xoá toàn bộ per-IP state (dùng giữa các kịch bản test)"),
    ], footnote="6 lớp model: bruteforce · dos · flood · legitimate · malformed · slowite")

    build_content_slide(prs, "Traffic Capture — traffic_capture.py", [
        (0, "Dùng tshark bắt traffic trên s1-eth11 (mirror port)"),
        (0, "BPF filter: tcp port 1883  ·  Display filter: mqtt"),
        (0, "Trích xuất 12 tshark fields, parse thành JSON, POST lên /predict"),
        (0, "4 aggregate features tính server-side trong ids_api.py:"),
        (1, "time_delta_mean / time_delta_std — thống kê inter-packet interval"),
        (1, "pkt_rate — số packet/giây trong cửa sổ 10s"),
        (1, "pub_to_conn_ratio — số PUBLISH / số CONNECT (phát hiện C2/exfil)"),
        (0, "Lợi thế mirror port: bắt traffic ở wire speed, không ảnh hưởng luồng chính"),
    ])

    build_content_slide(prs, "Ryu SDN Controller — ryu_controller.py", [
        (0, "L2 learning switch — OpenFlow 1.3"),
        (1, "Học MAC-to-port mapping, cài exact-match flow rule khi biết path"),
        (1, "Table-miss rule: gửi unmatched packet lên controller"),
        (0, "Port mirroring — s1-eth11:"),
        (1, "Copy ALL ingress traffic sang mirror port qua ovs-vsctl"),
        (1, "tshark bind tới eth11 — zero disruption to normal forwarding"),
        (0, "REST Flow Enforcer:"),
        (1, "POST /ids/block {ip} → cài DROP rule priority=200, permanent"),
        (1, "POST /ids/unblock {ip} → xóa DROP rule"),
        (1, "GET /ids/rules → liệt kê IP đang bị block"),
        (0, "Block ở data plane — drop tại OVS, không cần packet đến controller"),
    ])

    # ─── SECTION 4: Topology ─────────────────────────────────────────────────
    slide_section(prs, "04", "Topology Mininet",
                  "8 hosts · 1 OVS switch · 1 broker · 1 attacker")

    build_topology_slide(prs)

    build_content_slide(prs, "IP Addressing & Port Mapping", [
        (0, "Publishers: h1–h6  →  10.0.0.1 – 10.0.0.6  (eth1–eth6 on s1)"),
        (0, "Subscribers: h7–h8  →  10.0.0.7 – 10.0.0.8  (eth7–eth8 on s1)"),
        (0, "Broker: hbroker  →  10.0.0.10  (eth9 on s1)  ·  Mosquitto port 1883"),
        (0, "Attacker: hattacker  →  10.0.0.99  (eth10 on s1)"),
        (0, "Mirror port: s1-eth11  (không gắn host — tshark listen ở host OS)"),
        (0, "Controller: Ryu localhost:6633 (OpenFlow)  +  :8080 (REST)"),
        (0, "Bandwidth: broker & attacker link 100Mbps, host links 10Mbps"),
    ], footnote="Topology giống nhau trong Part2 và part3 — chỉ khác file topology.py chạy")

    # ─── SECTION 5: Attack Scenarios ─────────────────────────────────────────
    slide_section(prs, "05", "Kịch bản tấn công",
                  "Đơn lẻ · Kill-chain · Đa kẻ tấn công")

    build_attack_single_slide(prs)

    build_attack_killchain_slide(prs)

    build_attack_multi_slide(prs)

    # ─── SECTION 6: Evaluation ────────────────────────────────────────────────
    slide_section(prs, "06", "Kết quả đánh giá hệ thống",
                  "Per-class Precision/Recall/F1 · Latency")

    build_stats_slide(prs, "Tổng hợp kết quả đánh giá", [
        ("Accuracy",       "88.2%", "overall"),
        ("Macro F1",       "0.847",  ""),
        ("Macro Precision","0.853",  ""),
        ("Macro Recall",   "0.846",  ""),
        ("Macro FPR",      "2.4%",   "false positive rate"),
    ], note="Đánh giá trên 500 packets thực tế từ kịch bản Mininet  ·  Timestamp: 2026-04-22")

    build_content_slide(prs, "Kết quả theo lớp", [
        (0, "Normal (legitimate):  P=94.2%  R=94.2%  F1=94.2%  FPR=3.6%"),
        (0, "Flood (DoS):          P=90.2%  R=90.2%  F1=90.2%  FPR=1.9%"),
        (0, "C2 Malware:           P=88.6%  R=84.9%  F1=86.7%  FPR=1.9%"),
        (0, "Brute Force:          P=79.7%  R=90.8%  F1=84.9%  FPR=3.5%"),
        (0, "Port Scan:            P=79.6%  R=83.0%  F1=81.3%  FPR=2.2%"),
        (0, "Slow Drip:            P=79.4%  R=64.3%  F1=71.1%  FPR=1.5%"),
        (0, "→ Slow Drip thấp nhất vì traffic rất giống legitimate — đặc trưng low-and-slow"),
    ], footnote="Per-class breakdown từ part3/evaluation_results.json")

    build_image_slide(prs, "Per-class Performance",
                      IMG["per_class"],
                      caption="Precision / Recall / F1 theo từng lớp — flood & normal đạt tốt nhất, slow_drip recall thấp hơn",
                      img_left=Inches(1.0), img_top=Inches(1.55),
                      img_w=Inches(11.3))

    build_image_slide(prs, "Detection Latency",
                      IMG["latency"],
                      caption="Latency boxplot — thời gian từ khi gói tin đến khi IDS trả kết quả (bao gồm tshark parsing + model inference)",
                      img_left=Inches(2.0), img_top=Inches(1.55),
                      img_w=Inches(9.3))

    build_image_slide(prs, "Metrics Overview",
                      IMG["metrics_bar"],
                      caption="Precision / Recall / F1 tổng hợp — macro average across all classes",
                      img_left=Inches(1.5), img_top=Inches(1.55),
                      img_w=Inches(10))

    # ─── SECTION 7: Overhead Controller (Real End-to-End Benchmark) ─────────
    slide_section(prs, "07", "Overhead Controller",
                  subtitle="Benchmark thực — IDS API (Flask) + Ryu REST + E2E detect→block")
    build_two_image_slide(
        prs,
        title="Benchmark Throughput & Latency — Pipeline thực tế",
        img1=str(P2 / "bench_latency.png"),
        img2=str(P2 / "bench_throughput.png"),
        cap1="Latency p50/mean/p95 — 7 thành phần đo trên hệ thống đang chạy",
        cap2="Throughput — IDS /predict (4 luồng) + trần lý thuyết = 1000/mean",
    )
    build_content_slide(prs,
        title="Kết quả Benchmark — IDS API + Ryu Controller (chạy thật)",
        items=[
            (0, "Công cụ: benchmark_sdn.py — HTTP thật → ids_api.py:5000 + ryu_controller.py:8080"),
            (0, "Latency mean / p95 (200 mẫu, localhost)"),
            (1, "IDS /predict (normal):  5,21 / 6,30 ms"),
            (1, "IDS /predict (attack):  7,69 / 9,48 ms (+ vote + threat-score)"),
            (1, "Ryu /ids/block · /ids/unblock:  1,11 / 2,29 ms — FlowMod round-trip"),
            (1, "E2E detect→block:  5,57 / 6,05 ms"),
            (0, "Throughput (4 luồng, 5 s)"),
            (1, "IDS /predict: 188 req/s — 0 lỗi / 942 req"),
            (1, "Ryu REST trần ~900–1.160 req/s ⇒ controller không phải bottleneck"),
            (0, "Nhận xét"),
            (1, "Overhead SDN ≈ 0,4 ms; bottleneck là XGBoost (~5 ms), không phải OpenFlow"),
        ],
        footnote="Đo bằng benchmark_sdn.py — IDS_BLOCK_VOTES=10, IDS_BLOCK_CONF=0.8 (cấu hình demo Part2)"
    )

    # ─── SECTION 8: Conclusion ────────────────────────────────────────────────
    slide_section(prs, "08", "Kết luận & Hướng phát triển",
                  "")

    build_content_slide(prs, "Kết luận", [
        (0, "Đã xây dựng hệ thống IDS realtime cho mạng IoT/MQTT trên nền tảng SDN:"),
        (1, "Model XGBoost v5 — 16 features, 6 classes, test accuracy 83.2%, F1 83.8%"),
        (1, "IDS API với Threat Score — xử lý concurrent attack từ nhiều IP"),
        (1, "Ryu SDN controller — block tự động ở data plane, không cần human intervention"),
        (1, "Mininet topology — 10 nodes, MQTT broker, port mirroring"),
        (0, "3 kịch bản tấn công kiểm chứng:"),
        (1, "Đơn lẻ (5 loại) — kiểm tra từng classifier riêng"),
        (1, "Kill-chain (combined_attack.py) — kiểm tra vote/score dưới multi-label transition"),
        (1, "Multi-attacker — kiểm tra concurrent IP processing"),
        (0, "Hệ thống hoạt động end-to-end, block thành công trong kịch bản demo"),
    ])

    build_content_slide(prs, "Hướng phát triển", [
        (0, "Model:"),
        (1, "Thử LSTM / Transformer để khai thác temporal sequence của packets"),
        (1, "Online learning — cập nhật model khi gặp pattern mới (concept drift)"),
        (0, "Hệ thống:"),
        (1, "Triển khai multi-controller (ONOS/ODL) với failover"),
        (1, "Tích hợp TLS cho MQTT broker — decrypt trước khi IDS phân tích"),
        (1, "Dashboard realtime — visualize threat score, block events"),
        (0, "Đánh giá:"),
        (1, "Test trên traffic thực (không phải Mininet) — kiểm tra false positive rate"),
        (1, "Benchmark latency ở packet rate cao (>10k pkt/s)"),
        (1, "Adversarial testing — kẻ tấn công biết model và cố gắng né tránh"),
    ])

    # ─── Final ────────────────────────────────────────────────────────────────
    build_thank_you(prs)

    prs.save(str(OUT))
    print(f"✓ Saved → {OUT}")
    print(f"  Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
